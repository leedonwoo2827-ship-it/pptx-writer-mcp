# -*- coding: utf-8 -*-
"""
PPTX 생성 엔진

SlideSpec 리스트를 python-pptx로 렌더링합니다.
세션 관리: Presentation 객체를 메모리에 유지하며 슬라이드를 점진적으로 추가.
"""

import json
import os
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor


# ---------------------------------------------------------------------------
# 색상 맵 (기본)
# ---------------------------------------------------------------------------
DEFAULT_COLORS = {
    "red": "#dc2626",
    "green": "#16a34a",
    "blue": "#2563eb",
    "yellow": "#eab308",
    "black": "#000000",
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    """#RRGGBB → RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# PPTXGenerator 클래스
# ---------------------------------------------------------------------------

class PPTXGenerator:
    """PPTX 프레젠테이션 생성기."""

    def __init__(self, template_path: str = "", styles: dict = None, base_dir: str = ""):
        """
        Args:
            template_path: .pptx 템플릿 경로. 빈 문자열이면 빈 프레젠테이션
            styles: pptx-styles.json 내용
            base_dir: 이미지 등 상대경로 기준 디렉토리
        """
        if template_path and Path(template_path).exists():
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            # 와이드스크린 설정
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)

        self.styles = styles or {}
        self.base_dir = base_dir
        self.colors = {**DEFAULT_COLORS, **self.styles.get("colors", {})}
        self._layout_cache = {}
        self._content_area_cache = {}  # layout_name → {left, top, width, height}
        self._build_layout_cache()

    def _build_layout_cache(self):
        """레이아웃 이름 → SlideLayout 캐시."""
        mapping = self.styles.get("layout_mapping", {})
        for layout in self.prs.slide_layouts:
            name_lower = layout.name.lower().strip()
            self._layout_cache[name_lower] = layout
            # 매핑 역방향도 등록
            for key, val in mapping.items():
                if val.lower().strip() == name_lower:
                    self._layout_cache[key.lower()] = layout
            # 콘텐츠 영역 자동 감지
            self._detect_content_area(layout)

    def _detect_content_area(self, layout):
        """레이아웃의 콘텐츠 영역을 자동 감지.

        플레이스홀더가 없는 레이아웃에서:
        1. 콘텐츠 영역 사각형이 있으면 그것을 사용
        2. 없으면 헤더 영역 아래의 빈 공간을 계산
        """
        name = layout.name.lower().strip()

        # 스타일에서 명시적 content_area 설정 확인
        content_areas = self.styles.get("content_areas", {})
        if name in content_areas:
            self._content_area_cache[name] = content_areas[name]
            return

        # 플레이스홀더가 있으면 스킵 (플레이스홀더 기반 렌더링 사용)
        if list(layout.placeholders):
            return

        shapes = list(layout.shapes)
        if not shapes:
            return

        # 방법 1: 콘텐츠 영역 사각형 찾기 (top > 5cm이고 높이 > 10cm)
        best = None
        best_area = 0
        for shape in shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                continue
            w = shape.width or 0
            h = shape.height or 0
            t = shape.top or 0
            area = w * h
            if area > best_area and h > Emu(5000000) and t > Emu(2000000):
                best_area = area
                best = shape

        if best:
            self._content_area_cache[name] = {
                "left": best.left,
                "top": best.top,
                "width": best.width,
                "height": best.height,
            }
            return

        # 방법 2: 헤더 영역의 하단을 계산하여 콘텐츠 영역 추정
        header_bottom = 0
        footer_top = self.prs.slide_height
        for shape in shapes:
            bottom = (shape.top or 0) + (shape.height or 0)
            top = shape.top or 0
            # 상단 영역(top < 슬라이드 높이의 1/3)의 shape → 헤더
            if top < self.prs.slide_height // 3:
                header_bottom = max(header_bottom, bottom)
            # 하단 영역(top > 슬라이드 높이의 2/3)의 shape → 푸터
            if top > self.prs.slide_height * 2 // 3:
                footer_top = min(footer_top, top)

        if header_bottom > 0:
            margin = Emu(360000)  # 1cm 여백
            self._content_area_cache[name] = {
                "left": Emu(324000),    # 0.9cm
                "top": header_bottom + margin,
                "width": self.prs.slide_width - Emu(648000),  # 양쪽 0.9cm
                "height": footer_top - header_bottom - margin * 2,
            }

    def _get_content_area(self, layout_name: str):
        """레이아웃의 콘텐츠 영역 반환. 없으면 None."""
        name = layout_name.lower().strip()
        if name in self._content_area_cache:
            return self._content_area_cache[name]
        # 부분 매칭
        for cached_name, area in self._content_area_cache.items():
            if name in cached_name or cached_name in name:
                return area
        return None

    def _get_layout(self, layout_name: str):
        """레이아웃 이름으로 SlideLayout 조회. 폴백 체인 적용."""
        name = layout_name.lower().strip()

        # 1. 정확 매칭
        if name in self._layout_cache:
            return self._layout_cache[name]

        # 2. 부분 매칭
        for cached_name, layout in self._layout_cache.items():
            if name in cached_name or cached_name in name:
                return layout

        # 3. 인덱스 폴백
        layouts = self.prs.slide_layouts
        if name == "section_header" and len(layouts) > 2:
            return layouts[2]
        elif name == "blank" and len(layouts) > 6:
            return layouts[6]
        elif name == "title_slide" and len(layouts) > 0:
            return layouts[0]
        elif name == "two_content" and len(layouts) > 3:
            return layouts[3]
        elif name == "title_only" and len(layouts) > 5:
            return layouts[5]

        # 4. 기본: Title and Content (인덱스 1)
        if len(layouts) > 1:
            return layouts[1]
        return layouts[0]

    def add_slides(self, slide_specs: list) -> int:
        """SlideSpec 리스트를 프레젠테이션에 추가.

        Returns:
            추가된 슬라이드 수
        """
        count = 0
        for spec in slide_specs:
            layout_name = spec.get("layout", "title_content")
            layout = self._get_layout(layout_name)
            slide = self.prs.slides.add_slide(layout)

            # 실제 레이아웃 이름을 spec에 보존 (콘텐츠 영역 감지용)
            actual_layout_name = layout.name.lower().strip()
            spec["_actual_layout"] = actual_layout_name

            if layout_name in ("section_header",):
                self._render_section_header(slide, spec)
            elif layout_name in ("blank",):
                self._render_blank(slide, spec)
            else:
                self._render_content_slide(slide, spec)

            count += 1
        return count

    def save(self, output_path: str) -> str:
        """프레젠테이션을 파일로 저장."""
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)
        return output_path

    @property
    def slide_count(self) -> int:
        return len(self.prs.slides)

    def get_template_info(self) -> dict:
        """현재 템플릿의 레이아웃/플레이스홀더 정보."""
        info = {
            "slide_width": f"{self.prs.slide_width / 914400:.3f} inches",
            "slide_height": f"{self.prs.slide_height / 914400:.3f} inches",
            "layouts": [],
        }
        for i, layout in enumerate(self.prs.slide_layouts):
            layout_info = {
                "index": i,
                "name": layout.name,
                "placeholders": [],
            }
            for ph in layout.placeholders:
                layout_info["placeholders"].append({
                    "idx": ph.placeholder_format.idx,
                    "name": ph.name,
                    "width": f"{ph.width / 914400:.2f}in" if ph.width else "auto",
                    "height": f"{ph.height / 914400:.2f}in" if ph.height else "auto",
                })
            info["layouts"].append(layout_info)
        return info

    # -------------------------------------------------------------------
    # 렌더링: Section Header
    # -------------------------------------------------------------------
    def _render_section_header(self, slide, spec: dict):
        """Section Header 슬라이드 렌더링."""
        title = spec.get("title", "")
        layout_name = spec.get("layout", "section_header")

        # 플레이스홀더가 있으면 사용
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:  # Title
                ph.text = ""
                self._add_styled_text(
                    ph.text_frame, title,
                    self.styles.get("styles", {}).get("part_title", {}),
                )
                ph.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                break
        else:
            # 콘텐츠 영역 감지
            actual_layout = spec.get("_actual_layout", layout_name)
            content_area = self._get_content_area(actual_layout) or self._get_content_area(layout_name)
            if content_area:
                left = content_area["left"]
                top = content_area["top"]
                width = content_area["width"]
                height = Inches(2)
            else:
                left = Inches(1)
                top = Inches(2.5)
                width = self.prs.slide_width - Inches(2)
                height = Inches(2)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            self._add_styled_text(tf, title,
                                  self.styles.get("styles", {}).get("part_title", {}))
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------------
    # 렌더링: Content Slide (title_content)
    # -------------------------------------------------------------------
    def _render_content_slide(self, slide, spec: dict):
        """Title + Content 슬라이드 렌더링."""
        title = spec.get("title", "")
        subtitle = spec.get("subtitle", "")
        blocks = spec.get("blocks", [])
        layout_name = spec.get("layout", "title_content")

        # Title 플레이스홀더
        title_set = False
        content_ph = None
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            if idx == 0:  # Title
                display_title = title
                if subtitle:
                    display_title = f"{title} — {subtitle}" if title else subtitle
                ph.text = display_title
                self._style_placeholder_title(ph)
                title_set = True
            elif idx == 1:  # Content
                content_ph = ph

        # 콘텐츠 영역 감지 (실제 레이아웃 이름 우선)
        actual_layout = spec.get("_actual_layout", layout_name)
        content_area = self._get_content_area(actual_layout) or self._get_content_area(layout_name)

        if not title_set and title:
            display_title = f"{title} — {subtitle}" if subtitle else title
            if content_area:
                # 레이아웃의 헤더 영역에 타이틀 배치 (헤더 바 내부)
                title_left = content_area["left"]
                title_top = content_area["top"]
                title_width = content_area["width"]
                title_height = Emu(500000)  # ~1.4cm
            else:
                title_left = Inches(0.7)
                title_top = Inches(0.3)
                title_width = self.prs.slide_width - Inches(1.4)
                title_height = Inches(0.8)

            txBox = slide.shapes.add_textbox(
                title_left, title_top, title_width, title_height
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            self._add_styled_text(tf, display_title,
                                  self.styles.get("styles", {}).get("chapter_title", {}))

        # Content 렌더링
        if content_ph is not None:
            self._render_blocks_in_placeholder(content_ph, blocks)
        elif blocks:
            if content_area:
                # 레이아웃 콘텐츠 영역 기반 배치
                title_reserve = Emu(600000)  # 타이틀 공간 ~1.7cm
                self._render_blocks_direct(
                    slide, blocks,
                    top_offset=content_area["top"] + title_reserve,
                    margin_left=content_area["left"],
                    content_width=content_area["width"],
                    max_bottom=content_area["top"] + content_area["height"],
                )
            else:
                self._render_blocks_direct(slide, blocks, top_offset=Inches(1.4))

    # -------------------------------------------------------------------
    # 렌더링: Blank Slide
    # -------------------------------------------------------------------
    def _render_blank(self, slide, spec: dict):
        """Blank 슬라이드: 테이블/이미지 직접 배치."""
        title = spec.get("title", "")
        blocks = spec.get("blocks", [])
        layout_name = spec.get("layout", "blank")
        actual_layout = spec.get("_actual_layout", layout_name)

        content_area = self._get_content_area(actual_layout) or self._get_content_area(layout_name)

        if content_area:
            m_left = content_area["left"]
            m_width = content_area["width"]
            m_max = content_area["top"] + content_area["height"]
            top_offset = content_area["top"]
        else:
            m_left = Inches(0.7)
            m_width = self.prs.slide_width - Inches(1.4)
            m_max = None
            top_offset = Inches(0.5)

        if title:
            txBox = slide.shapes.add_textbox(
                m_left, top_offset, m_width, Inches(0.7)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            self._add_styled_text(tf, title,
                                  self.styles.get("styles", {}).get("chapter_title", {}))
            top_offset += Inches(0.8)

        self._render_blocks_direct(slide, blocks, top_offset=top_offset,
                                   margin_left=m_left, content_width=m_width,
                                   max_bottom=m_max)

    # -------------------------------------------------------------------
    # 블록 렌더링
    # -------------------------------------------------------------------
    def _render_blocks_in_placeholder(self, ph, blocks: list):
        """플레이스홀더의 TextFrame에 블록 렌더링."""
        tf = ph.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        first = True
        for block in blocks:
            block_type = block.get("type", "paragraph")

            if block_type == "bullet":
                if not first:
                    para = tf.add_paragraph()
                else:
                    para = tf.paragraphs[0]
                    first = False
                self._render_bullet_paragraph(para, block)

            elif block_type == "heading":
                if not first:
                    para = tf.add_paragraph()
                else:
                    para = tf.paragraphs[0]
                    first = False
                self._render_heading_paragraph(para, block)

            elif block_type == "paragraph":
                if not first:
                    para = tf.add_paragraph()
                else:
                    para = tf.paragraphs[0]
                    first = False
                self._render_text_paragraph(para, block)

            elif block_type in ("table", "image"):
                # 테이블/이미지는 플레이스홀더 내에서 처리 어려움
                # 텍스트 설명으로 대체 (나중에 직접 배치로 개선)
                if block_type == "table":
                    if not first:
                        para = tf.add_paragraph()
                    else:
                        para = tf.paragraphs[0]
                        first = False
                    self._render_table_as_text(para, block)
                # 이미지는 스킵 (직접 배치 필요)

    def _render_blocks_direct(self, slide, blocks: list, top_offset=None,
                              margin_left=None, content_width=None, max_bottom=None):
        """슬라이드에 블록을 직접 shape으로 배치."""
        if top_offset is None:
            top_offset = Inches(1.4)
        if margin_left is None:
            margin_left = Inches(0.7)
        if content_width is None:
            content_width = self.prs.slide_width - Inches(1.4)
        if max_bottom is None:
            max_bottom = self.prs.slide_height - Inches(0.5)

        current_top = top_offset

        for block in blocks:
            if current_top >= max_bottom:
                break

            block_type = block.get("type", "paragraph")

            if block_type == "table":
                height = self._add_table_shape(slide, block,
                                               margin_left, current_top, content_width)
                current_top += height + Inches(0.2)

            elif block_type == "image":
                height = self._add_image_shape(slide, block,
                                               margin_left, current_top, content_width)
                current_top += height + Inches(0.2)

            else:
                # 텍스트 블록 → 텍스트박스
                height = Inches(0.5)
                txBox = slide.shapes.add_textbox(
                    margin_left, current_top, content_width, height
                )
                tf = txBox.text_frame
                tf.word_wrap = True

                if block_type == "bullet":
                    self._render_bullet_paragraph(tf.paragraphs[0], block)
                elif block_type == "heading":
                    self._render_heading_paragraph(tf.paragraphs[0], block)
                else:
                    self._render_text_paragraph(tf.paragraphs[0], block)

                current_top += height + Inches(0.1)

    # -------------------------------------------------------------------
    # Paragraph 렌더링 헬퍼
    # -------------------------------------------------------------------
    def _render_bullet_paragraph(self, para, block: dict):
        """불릿 항목 렌더링."""
        level = block.get("level", 1)
        runs = block.get("runs", [])

        para.level = max(0, level - 1)
        style_key = f"bullet_{level}"
        style = self.styles.get("styles", {}).get(style_key,
                self.styles.get("styles", {}).get("bullet_1", {}))

        font_size = Pt(style.get("size", 14))
        space_before = Pt(style.get("space_before_pt", 4))
        space_after = Pt(style.get("space_after_pt", 2))

        para.space_before = space_before
        para.space_after = space_after

        # 불릿 기호
        bullet_chars = {1: "•", 2: "–", 3: "·"}
        prefix = bullet_chars.get(level, "•") + " "

        # 첫 번째 run에 불릿 기호 추가
        if runs:
            first_run = runs[0].copy()
            first_run["text"] = prefix + first_run.get("text", "")
            all_runs = [first_run] + runs[1:]
        else:
            all_runs = [{"text": prefix}]

        self._add_runs_to_paragraph(para, all_runs, style)

    def _render_heading_paragraph(self, para, block: dict):
        """헤딩(H3-H5) 렌더링."""
        level = block.get("level", 3)
        runs = block.get("runs", [])
        text = block.get("text", "")

        if not runs and text:
            runs = [{"text": text}]

        style = self.styles.get("styles", {}).get("section_title", {})
        sizes = {3: 18, 4: 16, 5: 14}
        override_style = {**style, "size": sizes.get(level, 16), "bold": True}

        para.space_before = Pt(8)
        para.space_after = Pt(4)

        self._add_runs_to_paragraph(para, runs, override_style)

    def _render_text_paragraph(self, para, block: dict):
        """일반 텍스트 렌더링."""
        runs = block.get("runs", [])
        style = self.styles.get("styles", {}).get("body", {})

        para.space_before = Pt(style.get("space_before_pt", 4))
        para.space_after = Pt(style.get("space_after_pt", 4))

        self._add_runs_to_paragraph(para, runs, style)

    def _render_table_as_text(self, para, block: dict):
        """테이블을 텍스트로 간이 렌더링 (플레이스홀더 내)."""
        headers = block.get("headers", [])
        rows = block.get("rows", [])
        text = " | ".join(headers) + "\n"
        for row in rows[:5]:  # 최대 5행 미리보기
            text += " | ".join(row) + "\n"
        if len(rows) > 5:
            text += f"... (총 {len(rows)}행)"

        run = para.add_run()
        run.text = text
        run.font.size = Pt(9)
        run.font.name = "맑은 고딕"

    def _add_runs_to_paragraph(self, para, runs: list, style: dict):
        """run 리스트를 paragraph에 추가."""
        font_name = style.get("font", "맑은 고딕")
        font_size = Pt(style.get("size", 12))
        default_bold = style.get("bold", False)
        default_color = style.get("color", "#333333")

        for i, run_data in enumerate(runs):
            if i == 0 and not para.runs:
                run = para.runs[0] if para.runs else para.add_run()
            else:
                run = para.add_run()

            run.text = run_data.get("text", "")
            run.font.name = font_name
            run.font.size = font_size
            run.font.bold = run_data.get("bold", default_bold)

            # 색상
            color_name = run_data.get("color", "")
            if color_name and color_name in self.colors:
                run.font.color.rgb = hex_to_rgb(self.colors[color_name])
            elif default_color:
                run.font.color.rgb = hex_to_rgb(default_color)

    def _add_styled_text(self, text_frame, text: str, style: dict):
        """TextFrame에 스타일이 적용된 텍스트 추가."""
        para = text_frame.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.name = style.get("font", "맑은 고딕")
        run.font.size = Pt(style.get("size", 28))
        run.font.bold = style.get("bold", False)
        if "color" in style:
            run.font.color.rgb = hex_to_rgb(style["color"])

    def _style_placeholder_title(self, ph):
        """Title 플레이스홀더 스타일."""
        style = self.styles.get("styles", {}).get("chapter_title", {})
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = style.get("font", "맑은 고딕")
                run.font.size = Pt(style.get("size", 28))
                run.font.bold = style.get("bold", True)
                if "color" in style:
                    run.font.color.rgb = hex_to_rgb(style["color"])

    # -------------------------------------------------------------------
    # Shape 렌더링: 테이블
    # -------------------------------------------------------------------
    def _add_table_shape(self, slide, block: dict,
                         left, top, width) -> int:
        """슬라이드에 테이블 shape 추가. 반환: 사용한 높이(EMU)."""
        headers = block.get("headers", [])
        rows = block.get("rows", [])
        header_runs = block.get("header_runs", [])
        row_runs = block.get("row_runs", [])

        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)

        if num_cols == 0:
            return 0

        row_height = Inches(0.35)
        table_height = row_height * num_rows
        max_height = self.prs.slide_height - top - Inches(0.5)
        table_height = min(table_height, max_height)

        tbl_shape = slide.shapes.add_table(
            num_rows, num_cols, left, top, width, table_height
        )
        table = tbl_shape.table

        # 열 폭 균등 분배
        col_width = int(width / num_cols)
        for col_idx in range(num_cols):
            table.columns[col_idx].width = col_width

        # 헤더 스타일
        header_style = self.styles.get("styles", {}).get("table_header", {})
        header_bg = hex_to_rgb(header_style.get("bg_color", "#1a1a2e"))
        header_font_color = hex_to_rgb(header_style.get("font_color", "#ffffff"))

        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = ""
            # 배경색
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg

            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER

            if header_runs and col_idx < len(header_runs):
                self._add_cell_runs(para, header_runs[col_idx],
                                    header_style, header_font_color)
            else:
                run = para.add_run()
                run.text = header_text
                run.font.name = header_style.get("font", "맑은 고딕")
                run.font.size = Pt(header_style.get("size", 10))
                run.font.bold = True
                run.font.color.rgb = header_font_color

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 데이터 행
        body_style = self.styles.get("styles", {}).get("table_body", {})
        alt_color = body_style.get("alt_row_color", "#f0f4f8")

        for row_idx, row_cells in enumerate(rows):
            for col_idx, cell_text in enumerate(row_cells):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = ""

                # 교대행 색상
                if row_idx % 2 == 1 and alt_color:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = hex_to_rgb(alt_color)

                para = cell.text_frame.paragraphs[0]
                align_str = body_style.get("align", "center")
                para.alignment = {
                    "center": PP_ALIGN.CENTER,
                    "left": PP_ALIGN.LEFT,
                    "right": PP_ALIGN.RIGHT,
                }.get(align_str, PP_ALIGN.CENTER)

                if row_runs and row_idx < len(row_runs) and col_idx < len(row_runs[row_idx]):
                    self._add_cell_runs(para, row_runs[row_idx][col_idx],
                                        body_style)
                else:
                    run = para.add_run()
                    run.text = cell_text
                    run.font.name = body_style.get("font", "맑은 고딕")
                    run.font.size = Pt(body_style.get("size", 9))
                    if "color" in body_style:
                        run.font.color.rgb = hex_to_rgb(body_style["color"])

                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        return int(table_height)

    def _add_cell_runs(self, para, runs: list, style: dict,
                       override_color: RGBColor = None):
        """셀에 색상 run 추가."""
        for run_data in runs:
            run = para.add_run()
            run.text = run_data.get("text", "")
            run.font.name = style.get("font", "맑은 고딕")
            run.font.size = Pt(style.get("size", 10))
            run.font.bold = run_data.get("bold", style.get("bold", False))

            color_name = run_data.get("color", "")
            if color_name and color_name in self.colors:
                run.font.color.rgb = hex_to_rgb(self.colors[color_name])
            elif override_color:
                run.font.color.rgb = override_color
            elif "color" in style:
                run.font.color.rgb = hex_to_rgb(style["color"])

    # -------------------------------------------------------------------
    # Shape 렌더링: 이미지
    # -------------------------------------------------------------------
    def _add_image_shape(self, slide, block: dict,
                         left, top, width) -> int:
        """슬라이드에 이미지 shape 추가. 반환: 사용한 높이(EMU)."""
        img_path = block.get("path", "")
        caption = block.get("caption", "")

        # 이미지 경로 해석
        resolved = self._resolve_image_path(img_path)
        if not resolved or not Path(resolved).exists():
            # 이미지 없으면 캡션만 텍스트박스로
            height = Inches(0.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            run = tf.paragraphs[0].add_run()
            run.text = f"[이미지: {caption or img_path}]"
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
            return int(height)

        # 이미지 삽입
        max_height = self.prs.slide_height - top - Inches(1)
        max_width = width

        try:
            pic = slide.shapes.add_picture(
                str(resolved), left, top, max_width
            )
            # 비율 유지
            ratio = pic.image.size[0] / pic.image.size[1]
            if pic.height > max_height:
                pic.height = int(max_height)
                pic.width = int(max_height * ratio)
            # 가운데 정렬
            pic.left = int(left + (width - pic.width) / 2)

            img_height = pic.height

            # 캡션
            if caption:
                cap_top = top + img_height + Inches(0.1)
                txBox = slide.shapes.add_textbox(left, cap_top, width, Inches(0.4))
                tf = txBox.text_frame
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                run = tf.paragraphs[0].add_run()
                run.text = caption
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
                img_height += Inches(0.5)

            return int(img_height)

        except Exception:
            height = Inches(0.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            run = tf.paragraphs[0].add_run()
            run.text = f"[이미지 로드 실패: {img_path}]"
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0xcc, 0x00, 0x00)
            return int(height)

    def _resolve_image_path(self, img_path: str) -> Optional[str]:
        """이미지 경로 해석: 절대경로 → base_dir 상대 → 없으면 None."""
        p = Path(img_path)
        if p.is_absolute() and p.exists():
            return str(p)

        if self.base_dir:
            candidate = Path(self.base_dir) / img_path
            if candidate.exists():
                return str(candidate)

        return None
